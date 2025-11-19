import pptx
import pptx.dml
import pptx.dml.effect
import pptx.enum
import pptx.enum.dml
import pptx.enum.shapes
import pptx.enum.text
import pptx.slide
import math

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement


class Component:
    def __init__(self):
        pass

    def render(self, slide: pptx.slide.Slide):
        raise NotImplementedError


class FLAG(Component):
    def __init__(self, name: str, color: RGBColor, flag: bool, left: Cm, top: Cm):
        super().__init__()
        self.name = name
        self.color = color
        self.left = left
        self.top = top
        self.flag = flag

        self.width = Cm(3)
        self.height = Cm(1)

    def set(self, flag: bool):
        self.flag = flag

    def render(self, slide: pptx.slide.Slide):
        shapes = slide.shapes

        shapes_to_group = []

        # 1. 繪製背景長方形 (底色為 color)
        rect = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            self.left,
            self.top,
            self.width,
            self.height,
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = self.color
        rect.line.color.rgb = RGBColor(0, 0, 0)  # 黑框
        rect.line.width = Pt(1)
        rect.shadow.inherit = False
        shapes_to_group.append(rect)

        # 排版參數
        margin_x = Cm(0.2)  # 左右邊距

        # 2. 繪製指示燈 (圓形) - 置右
        # 直徑設為高度的 60%
        diameter = self.height * 0.6
        circle_right_margin = margin_x

        circle_left = self.left + self.width - diameter - circle_right_margin
        circle_top = self.top + (self.height - diameter) / 2  # 垂直置中

        indicator = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, circle_left, circle_top, diameter, diameter
        )
        indicator.line.color.rgb = RGBColor(0, 0, 0)  # 黑框
        indicator.line.width = Pt(1.5)
        indicator.fill.solid()
        indicator.shadow.inherit = False
        shapes_to_group.append(indicator)

        # 根據 flag 決定顏色
        if self.flag:
            indicator.fill.fore_color.rgb = RGBColor(0, 200, 0)  # 亮綠色
        else:
            indicator.fill.fore_color.rgb = RGBColor(50, 50, 50)  # 黑色

        # 3. 繪製名稱 (文字方塊) - 置左
        # 文字可用寬度 = 總寬 - 指示燈寬 - 邊距
        text_width = self.width - diameter - (margin_x * 3)

        tb = shapes.add_textbox(self.left + margin_x, self.top, text_width, self.height)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.word_wrap = False
        shapes_to_group.append(tb)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.text = self.name
        p.font.name = "Tahoma"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # 5. 執行群組化操作
        self.group_shapes(slide, shapes_to_group)

    def group_shapes(self, slide, shapes_to_group):
        """
        使用 XML 操作將指定的 shapes 列表組合成一個 PowerPoint Group Shape。
        注意：這是 python-pptx 的非官方 workaround。
        """
        if not shapes_to_group:
            return

        # 取得投影片的形狀樹 (spTree)
        spTree = slide.shapes._spTree

        # 建立群組元素 <p:grpSp>
        grpSp = OxmlElement("p:grpSp")

        # 設定群組的非視覺屬性 (nvGrpSpPr)
        nvGrpSpPr = OxmlElement("p:nvGrpSpPr")
        cNvPr = OxmlElement("p:cNvPr")
        cNvPr.set("id", "1")  # ID 隨意，PPT 會自動修正
        cNvPr.set("name", f"{self.name}_Group")  # 設定群組名稱
        nvGrpSpPr.append(cNvPr)
        nvGrpSpPr.append(OxmlElement("p:cNvGrpSpPr"))
        nvGrpSpPr.append(OxmlElement("p:nvPr"))
        grpSp.append(nvGrpSpPr)

        # 設定群組屬性 (grpSpPr) 與座標轉換 (xfrm)
        grpSpPr = OxmlElement("p:grpSpPr")
        xfrm = OxmlElement("a:xfrm")

        # 設定群組座標原點與大小 (這裡設為 0，並透過 chOff/chExt 對應，讓子形狀維持絕對座標)
        for tag in ["a:off", "a:ext", "a:chOff", "a:chExt"]:
            elem = OxmlElement(tag)
            elem.set("x" if "off" in tag.lower() else "cx", "0")
            elem.set("y" if "off" in tag.lower() else "cy", "0")
            xfrm.append(elem)

        grpSpPr.append(xfrm)
        grpSp.append(grpSpPr)

        # 將形狀從 spTree 移動到 grpSp
        for shape in shapes_to_group:
            spTree.remove(shape.element)  # 從投影片移除
            grpSp.append(shape.element)  # 加入群組

        # 將群組加入投影片
        spTree.append(grpSp)


class PE(Component):
    def __init__(
        self, name: str, left: Cm, top: Cm, width: Cm, height: Cm, PE_num: int
    ):
        super().__init__()
        self.name = name
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.PE_num = PE_num

        self.data = [[] for _ in range(PE_num)]
        self.cnt = [0 for _ in range(PE_num)]

    def write(self, text: str, cnt: int, color: RGBColor = None):
        """
        將資料 (text, color) 寫入指定索引的 PE 單元
        """
        if self.ready():
            for index in range(self.PE_num):
                if self.cnt[index] == 0:
                    self.cnt[index] = cnt
                    self.data[index].append((text, color))
                    break
        else:
            print(f"PE '{self.name}' 仍有工作未完成")

    def count(self):
        self.cnt = [x - 1 if x > 0 else x for x in self.cnt]

    def ready(self):
        return any(x == 0 for x in self.cnt)

    def render(self, slide: pptx.slide.Slide):
        shapes = slide.shapes

        # 用來收集需要群組的形狀
        shapes_to_group = []  # 背景結構 (外框、標題、圓形)
        shapes_to_group_ = []  # 計數器 (紅色數字)

        # 用來收集不需要群組的形狀 (Data 文字)，確保它們在最上層
        data_shapes = []

        # 1. 繪製最外層的圓角黑框長方形
        outer_box = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            self.left,
            self.top,
            self.width,
            self.height,
        )
        outer_box.fill.background()
        outer_box.line.color.rgb = RGBColor(0, 0, 0)
        outer_box.line.width = Pt(1.5)
        shapes_to_group.append(outer_box)

        # 2. 放置標題
        title_height = Cm(0.7)
        tb_title = shapes.add_textbox(
            self.left, self.top + Cm(0.1), self.width, title_height
        )
        tf_title = tb_title.text_frame
        tf_title.margin_bottom = Cm(0)
        tf_title.margin_top = Cm(0)
        tf_title.margin_left = Cm(0)
        tf_title.margin_right = Cm(0)
        tf_title.word_wrap = False
        p_title = tf_title.paragraphs[0]
        p_title.alignment = PP_ALIGN.CENTER
        p_title.text = self.name
        p_title.font.name = "Tahoma"
        p_title.font.size = Pt(10)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(0, 0, 0)
        shapes_to_group.append(tb_title)

        # 3. 計算圓形網格佈局
        if self.PE_num > 0:
            num_cols = int(math.ceil(math.sqrt(self.PE_num)))
            num_rows = int(math.ceil(self.PE_num / num_cols))

            area_left = self.left
            area_top = self.top + title_height
            area_width = self.width
            area_height = self.height - title_height - Cm(0.1)

            cell_size = min(area_width / num_cols, area_height / num_rows)
            circle_dia = cell_size * 0.8
            cell_padding = (cell_size - circle_dia) / 2

            grid_width = num_cols * cell_size
            grid_height = num_rows * cell_size
            grid_offset_x = (area_width - grid_width) / 2
            grid_offset_y = (area_height - grid_height) / 2

            grid_start_x = area_left + grid_offset_x
            grid_start_y = area_top + grid_offset_y

            # 4. 繪製 PE 單元
            pe_index = 0
            for r in range(num_rows):
                for c in range(num_cols):
                    if pe_index >= self.PE_num:
                        break

                    circle_left = grid_start_x + (c * cell_size) + cell_padding
                    circle_top = grid_start_y + (r * cell_size) + cell_padding

                    # 4.1 圓形 (加入群組)
                    circle_shape = shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.OVAL,
                        circle_left,
                        circle_top,
                        circle_dia,
                        circle_dia,
                    )
                    circle_shape.fill.background()
                    circle_shape.line.color.rgb = RGBColor(0, 0, 0)
                    circle_shape.line.width = Pt(1)
                    shapes_to_group.append(circle_shape)

                    # 4.2 Data 文字 (處理多筆資料)
                    data_list = self.data[pe_index]  # 這現在是一個 List

                    # 過濾掉空的資料 (如果有的話)
                    valid_items = [d for d in data_list if d[0]]

                    if valid_items:
                        tb = shapes.add_textbox(
                            circle_left, circle_top, circle_dia, circle_dia
                        )
                        tf = tb.text_frame
                        tf.margin_bottom = Cm(0)
                        tf.margin_top = Cm(0)
                        tf.margin_left = Cm(0)
                        tf.margin_right = Cm(0)
                        tf.word_wrap = False
                        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

                        # 處理第一筆資料 (使用預設的第一個段落，避免空行)
                        first_text, first_color = valid_items[0]
                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        p.text = str(first_text)
                        p.font.name = "Tahoma"
                        p.font.size = Pt(12)
                        p.font.bold = True
                        if first_color:
                            p.font.color.rgb = first_color

                        # 處理後續資料 (新增段落)
                        for text, color in valid_items[1:]:
                            p = tf.add_paragraph()
                            p.alignment = PP_ALIGN.CENTER
                            p.text = str(text)
                            p.font.name = "Tahoma"
                            p.font.size = Pt(12)
                            p.font.bold = True
                            if color:
                                p.font.color.rgb = color

                        # 關鍵：收集到 data_shapes
                        data_shapes.append(tb)

                    # 4.3 計數器 (加入第二群組)
                    cnt_val = self.cnt[pe_index]
                    if cnt_val >= 0:
                        cnt_box_size = circle_dia * 0.5
                        cnt_left = circle_left + (circle_dia * 0.5)
                        cnt_top = circle_top - (circle_dia * 0.1)

                        tb_cnt = shapes.add_textbox(
                            cnt_left, cnt_top, cnt_box_size, cnt_box_size
                        )
                        tf_cnt = tb_cnt.text_frame
                        tf_cnt.margin_right = Cm(0.1)
                        tf_cnt.word_wrap = False
                        p_cnt = tf_cnt.paragraphs[0]
                        p_cnt.alignment = PP_ALIGN.RIGHT
                        p_cnt.text = str(cnt_val)
                        p_cnt.font.name = "Tahoma"
                        p_cnt.font.size = Pt(10)
                        p_cnt.font.bold = True
                        p_cnt.font.color.rgb = (
                            RGBColor(255, 0, 0) if cnt_val > 0 else RGBColor(0, 0, 0)
                        )

                        shapes_to_group_.append(tb_cnt)

                    pe_index += 1
                if pe_index >= self.PE_num:
                    break

        # 5. 執行群組化操作
        self.group_shapes(slide, shapes_to_group)
        self.group_shapes(slide, shapes_to_group_)

        # 6. 修正圖層順序 (Z-Order)
        self.bring_to_front(slide, data_shapes)

    def group_shapes(self, slide, shapes_to_group):
        """
        使用 XML 操作將指定的 shapes 列表組合成一個 PowerPoint Group Shape。
        """
        if not shapes_to_group:
            return

        spTree = slide.shapes._spTree
        grpSp = OxmlElement("p:grpSp")

        nvGrpSpPr = OxmlElement("p:nvGrpSpPr")
        cNvPr = OxmlElement("p:cNvPr")
        cNvPr.set("id", "1")
        cNvPr.set("name", f"{self.name}_Group")
        nvGrpSpPr.append(cNvPr)
        nvGrpSpPr.append(OxmlElement("p:cNvGrpSpPr"))
        nvGrpSpPr.append(OxmlElement("p:nvPr"))
        grpSp.append(nvGrpSpPr)

        grpSpPr = OxmlElement("p:grpSpPr")
        xfrm = OxmlElement("a:xfrm")

        for tag in ["a:off", "a:ext", "a:chOff", "a:chExt"]:
            elem = OxmlElement(tag)
            elem.set("x" if "off" in tag.lower() else "cx", "0")
            elem.set("y" if "off" in tag.lower() else "cy", "0")
            xfrm.append(elem)

        grpSpPr.append(xfrm)
        grpSp.append(grpSpPr)

        for shape in shapes_to_group:
            spTree.remove(shape.element)
            grpSp.append(shape.element)

        spTree.append(grpSp)

    def bring_to_front(self, slide, shapes):
        """
        將指定的形狀移到最上層
        """
        spTree = slide.shapes._spTree
        for shape in shapes:
            spTree.remove(shape.element)
            spTree.append(shape.element)


class SRAM(Component):
    def __init__(
        self,
        row: int,
        col: int,
        # --- 修改這裡 (型別提示) ---
        left: Cm,
        top: Cm,
        width: Cm,
        height: Cm,
        interleave: bool = False,
        bottom_start: bool = True,
    ):
        self.row = row
        self.col = col
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.interleave = interleave  # 是否啟用帶狀欄
        self.bottom_start = bottom_start  # 座標是否從底部開始計算

        # 將 data store 改為儲存 (text, color) 的元組 (tuple)
        self.data = [
            [("", RGBColor(0, 0, 0), "") for _ in range(col)] for _ in range(row)
        ]

    def get_position(self, i: int, j: int):
        """
        獲取 (i, j) 儲存格的左上角 (x, y) 座標。
        i, j 皆為 0-based 索引 (從 top-left 開始)。
        - i: row 索引 (0 to self.row - 1)
        - j: col 索引 (0 to self.col - 1)

        會根據 self.bottom_start 決定 y 座標的對應方式。
        """
        x = self.left + j * self.width / self.col

        cell_height = self.height / self.row
        if self.bottom_start:
            # bottom-up 映射: data[0] (頂部資料) 對應到
            # 視覺上的最後一列 (visual_row_index = self.row - 1)
            visual_row_index = (self.row - 1) - i
            y = self.top + (visual_row_index * cell_height)
        else:
            # top-down 映射: data[0] 對應到 視覺上的第一列 (i=0)
            y = self.top + i * cell_height

        return x, y

    def render(self, slide: pptx.slide.Slide):
        """
        將这个 SRAM 表格新增到指定的投影片上
        """
        shapes = slide.shapes
        table = shapes.add_table(
            self.row,
            self.col,
            self.left,
            self.top,
            self.width,
            self.height,
        ).table

        table.first_row = False

        # 根據 self.interleave 啟用帶狀欄
        if self.interleave:
            table.vert_banding = True  # 啟用帶狀欄 (Banded Columns)
            table.horz_banding = False
        else:
            table.vert_banding = False
            table.horz_banding = False

        col_width = self.width / self.col
        for i in range(self.col):
            table.columns[i].width = int(col_width)

        row_height = self.height / self.row
        for i in range(self.row):
            table.rows[i].height = int(row_height)

        for i in range(self.row):
            for j in range(self.col):
                data, color, state = self.data[i][j]
                ii = i
                if self.bottom_start:
                    ii = self.row - 1 - i
                if state == "w":
                    table.cell(ii, j).fill.solid()
                    table.cell(ii, j).fill.fore_color.rgb = RGBColor(255, 0, 0)
                    self.data[i][j] = data, color, ""
                elif state == "r":
                    table.cell(ii, j).fill.solid()
                    table.cell(ii, j).fill.fore_color.rgb = RGBColor(0, 255, 0)
                    self.data[i][j] = data, color, ""

        self.render_data(slide)

    def write(self, i: int, j: int, text: str, color: RGBColor = None):
        """
        將資料 (text, color) 寫入內部的 self.data 暫存
        i, j 為 0-based 索引 (從 top-left 開始)
        """
        if 0 <= i < self.row and 0 <= j < self.col:
            self.data[i][j] = (text, color, "w")
        else:
            print(f"SRAM write 錯誤: 索引 ({i}, {j}) 超出範圍")

    def read(self, i: int, j: int):
        """
        將資料 (text, color) 寫入內部的 self.data 暫存
        i, j 為 0-based 索引 (從 top-left 開始)
        """
        if 0 <= i < self.row and 0 <= j < self.col:
            read_data, read_color, _ = self.data[i][j]
            self.data[i][j] = ("", RGBColor(0, 0, 0), "r")
            return read_data, read_color
        else:
            print(f"SRAM read 錯誤: 索引 ({i}, {j}) 超出範圍")
            return None, None

    def render_data(self, slide: pptx.slide.Slide):
        """
        遍歷 self.data，將所有資料繪製為文字方塊
        """
        if slide is None:
            print("SRAM 錯誤: 必須先呼叫 .add(slide) 才能 render_data")
            return

        cell_width = self.width / self.col
        cell_height = self.height / self.row

        for i in range(self.row):
            for j in range(self.col):
                text, color, _ = self.data[i][j]

                # 如果 data 不是 None 或空字串
                if text:
                    x, y = self.get_position(i, j)

                    if j % 2:
                        y += 0.5 * cell_height

                    tb = slide.shapes.add_textbox(x, y, cell_width, cell_height)
                    tf = tb.text_frame

                    # 清除邊界並置中
                    tf.margin_bottom = Cm(0)
                    tf.margin_top = Cm(0)
                    tf.margin_left = Cm(0)
                    tf.margin_right = Cm(0)
                    tf.word_wrap = False  # 避免自動換行

                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER  # 水平置中
                    # (垂直置中是 tf.vertical_anchor)

                    p.text = text
                    p.font.name = "Tahoma"  # 設定字體
                    p.font.size = Pt(12)  # 設定字體大小
                    p.font.bold = True

                    if color:
                        p.font.color.rgb = color
