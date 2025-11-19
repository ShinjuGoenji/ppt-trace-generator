from pptx import Presentation
import pptx
import pptx.presentation
from pptx.util import Cm

from pptx.dml.color import RGBColor
from componets import SRAM, PE, FLAG
from util import COLORS


# --- parameters --- #
OUTPUT_FILE = "generated_trace.pptx"
TITLE = "CRT Trace"

# Componenets
sram = SRAM(
    row=7,
    col=32,
    left=Cm(6),
    top=Cm(3),
    width=Cm(21),
    height=Cm(7),
    interleave=True,
    ptr={
        "mod_small_unsigned_start_ptr": {"i": 0, "j": 0, "color": COLORS["YELLOW"]},
        "mod_small_unsigned_end_ptr": {"i": 0, "j": 0, "color": COLORS["ORANGE"]},
        "add_mul_small_start_ptr": {"i": 1, "j": 0, "color": COLORS["GREEN"]},
        "add_mul_small_end_ptr": {"i": 1, "j": 0, "color": COLORS["CYAN"]},
    },
)

add_mul_small = PE(
    name="ADD_MUL_SMALL",
    left=Cm(12),
    top=Cm(12),
    width=Cm(4),
    height=Cm(4),
    PE_num=4,
)

mod_small_unsigned = PE(
    name="MOD_SMALL_UNSIGNED",
    left=Cm(18),
    top=Cm(12),
    width=Cm(4),
    height=Cm(4),
    PE_num=4,
)

r_flag_0 = FLAG(
    name="Read", color=RGBColor(178, 178, 190), flag=False, left=Cm(2), top=Cm(3)
)

w_flag_0 = FLAG(
    name="Write", color=RGBColor(178, 178, 190), flag=False, left=Cm(2), top=Cm(4.5)
)

r_flag_1 = FLAG(
    name="Read", color=RGBColor(227, 229, 237), flag=False, left=Cm(2), top=Cm(6)
)

w_flag_1 = FLAG(
    name="Write", color=RGBColor(227, 229, 237), flag=False, left=Cm(2), top=Cm(7.5)
)

write_buffer = SRAM(
    row=6,
    col=2,
    left=Cm(6.5),
    top=Cm(12),
    width=Cm(4),
    height=Cm(1),
    interleave=True,
    bottom_start=False,
)

add_mul_small_buffer = SRAM(
    row=4,
    col=1,
    left=Cm(4.5),
    top=Cm(12),
    width=Cm(1),
    height=Cm(4),
    interleave=False,
    bottom_start=False,
)

components = [
    sram,
    add_mul_small,
    mod_small_unsigned,
    r_flag_0,
    r_flag_1,
    w_flag_0,
    w_flag_1,
    write_buffer,
    add_mul_small_buffer,
]


# --- MAIN LOGIC --- #
def main_logic(prs: pptx.presentation.Presentation):
    inputs = {}
    for i in range(32 * 4):
        inputs[i] = (f"x{int(i/32)}_{i%32}", COLORS["BLUE"])
    input_cnt = 0

    write_buffer_ptr = [0, 0]
    add_mul_small_valid = [False, False, False, False]

    MAX_CYCLE = 70
    for cycle in range(MAX_CYCLE):

        input_data, input_color = inputs[input_cnt]
        u = int(input_cnt / 32)
        v = input_cnt % 32

        w_flag_0.flag = False
        w_flag_1.flag = False
        r_flag_0.flag = False
        r_flag_1.flag = False

        for k, d in enumerate(add_mul_small.data):
            if len(d) == 0:
                continue

            text, color = d[0]

            if text == "":
                continue

            i = int(text.replace("x", "").split("_")[0])
            j = int(text.split("_")[1])
            cnt = add_mul_small.cnt[k]

            if cnt > i + 1:
                add_mul_small.cnt[k] -= 1
            elif cnt == 1:
                add_mul_small.cnt[k] -= 1
            elif add_mul_small_valid[k] and cnt != 0:
                add_mul_small.cnt[k] -= 1
                add_mul_small_valid[k] = False

        mod_small_unsigned.count()

        sram_write = []
        mod_small_unsigned_write = []
        add_mul_small_write = []
        write_buffer_write = []
        add_mul_small_buffer_write = []

        # inputs
        if sram.ptr_value(
            "mod_small_unsigned_start_ptr", 32
        ) < input_cnt and sram.ptr_value(
            "mod_small_unsigned_start_ptr", 32
        ) < sram.ptr_value(
            "add_mul_small_end_ptr", 32
        ):
            if mod_small_unsigned.ready():
                read_data, read_color = sram.read(
                    sram.ptr["mod_small_unsigned_start_ptr"]["i"],
                    sram.ptr["mod_small_unsigned_start_ptr"]["j"],
                )
                mod_small_unsigned_write.append(
                    (
                        read_data,
                        (
                            1
                            + 3 * (sram.ptr["mod_small_unsigned_start_ptr"]["i"] + 1)
                            + 1
                        ),
                        read_color,
                    )
                )
                if sram.ptr_value("mod_small_unsigned_start_ptr", 32) % 2 == 0:
                    r_flag_0.flag = True
                else:
                    r_flag_1.flag = True
        elif sram.ptr_value("add_mul_small_start_ptr", 32) < input_cnt:
            if add_mul_small.ready():
                if (
                    sram.ptr_value("add_mul_small_start_ptr", 32) % 2 == 0
                    and not r_flag_0.flag
                ):
                    read_data, read_color = sram.read(
                        sram.ptr["add_mul_small_start_ptr"]["i"],
                        sram.ptr["add_mul_small_start_ptr"]["j"],
                    )
                    add_mul_small_write.append(
                        (
                            read_data,
                            (
                                1
                                + 3 * (sram.ptr["add_mul_small_start_ptr"]["i"] + 1)
                                + 1
                            ),
                            read_color,
                        )
                    )
                    r_flag_0.flag = True
                elif (
                    sram.ptr_value("add_mul_small_start_ptr", 32) % 2 == 1
                    and not r_flag_1.flag
                ):
                    read_data, read_color = sram.read(
                        sram.ptr["add_mul_small_start_ptr"]["i"],
                        sram.ptr["add_mul_small_start_ptr"]["j"],
                    )
                    add_mul_small_write.append(
                        (
                            read_data,
                            (
                                1
                                + 3 * (sram.ptr["add_mul_small_start_ptr"]["i"] + 1)
                                + 1
                            ),
                            read_color,
                        )
                    )
                    r_flag_1.flag = True

        if input_cnt < len(inputs):
            if u == 0:
                if (
                    mod_small_unsigned.ready()
                    and sram.ptr_value("mod_small_unsigned_start_ptr", 32) == input_cnt
                ):
                    mod_small_unsigned_write.append(
                        (
                            input_data,
                            5,
                            input_color,
                        )
                    )
                else:
                    sram_write.append((u, v, input_data, input_color))
                    if v % 2 == 0:
                        w_flag_0.flag = True
                    else:
                        w_flag_1.flag = True
            else:
                if (
                    add_mul_small.ready()
                    and sram.ptr_value("add_mul_small_start_ptr", 32) == input_cnt
                ):
                    add_mul_small_write.append((input_data, (3 + u + 1), input_color))
                else:
                    sram_write.append((u, v, input_data, input_color))
                    if v % 2 == 0:
                        w_flag_0.flag = True
                    else:
                        w_flag_1.flag = True

        for k, v in enumerate(add_mul_small.data):
            if len(v) == 0:
                continue

            text, color = v[0]
            if text == "":
                continue

            i = int(text.replace("x", "").split("_")[0])
            j = int(text.split("_")[1])
            cnt = add_mul_small.cnt[k]
            target_cnt = i + 1 if cnt > i + 1 else cnt

            if (
                not add_mul_small_valid[k]
                and sram.ptr_value("mod_small_unsigned_start_ptr", 32)
                > (i - target_cnt + 1) * 32 + j
            ):
                if j % 2 == 0 and not r_flag_0.flag:
                    data, color = sram.read(i - target_cnt + 1, j)
                    r_flag_0.flag = True
                    add_mul_small_valid[k] = True
                    if cnt > i + 1:
                        add_mul_small_buffer_write.append((k, 0, data, color))
                    else:
                        add_mul_small.write(data, cnt, color)
                elif j % 2 == 1 and not r_flag_1.flag:
                    data, color = sram.read(i - target_cnt + 1, j)
                    r_flag_1.flag = True
                    add_mul_small_valid[k] = True
                    if cnt > i + 1:
                        add_mul_small_buffer_write.append((k, 0, data, color))
                    else:
                        add_mul_small.write(data, cnt, color)

        if w_flag_0.flag == False and write_buffer_ptr[0] > 0:
            for k in range(write_buffer_ptr[0]):
                text, color = write_buffer.read(k, 0)
                if k == 0:
                    i = int(text.replace("x", "").split("_")[0])
                    j = int(text.split("_")[1])
                    sram_write.append((i, j, text, color))
                    w_flag_0.flag = True

                    if color == COLORS["YELLOW"]:
                        if sram.ptr["mod_small_unsigned_end_ptr"]["j"] + 1 >= 32:
                            sram.ptr["mod_small_unsigned_end_ptr"]["i"] += 1
                            sram.ptr["mod_small_unsigned_end_ptr"]["j"] = 0
                        else:
                            sram.ptr["mod_small_unsigned_end_ptr"]["j"] += 1
                    elif (
                        color == COLORS["GREEN"]
                        and i == sram.ptr["add_mul_small_end_ptr"]["i"]
                    ):
                        if sram.ptr["add_mul_small_end_ptr"]["j"] + 1 >= 32:
                            sram.ptr["add_mul_small_end_ptr"]["i"] += 1
                            sram.ptr["add_mul_small_end_ptr"]["j"] = 0
                        else:
                            sram.ptr["add_mul_small_end_ptr"]["j"] += 1
                else:
                    write_buffer_write.append((k - 1, 0, text, color))
            write_buffer_ptr[0] -= 1
        if w_flag_1.flag == False and write_buffer_ptr[1] > 0:
            for k in range(write_buffer_ptr[1]):
                text, color = write_buffer.read(k, 1)
                if k == 0:
                    i = int(text.replace("x", "").split("_")[0])
                    j = int(text.split("_")[1])
                    sram_write.append((i, j, text, color))
                    w_flag_1.flag = True

                    if color == COLORS["YELLOW"]:
                        if sram.ptr["mod_small_unsigned_end_ptr"]["j"] + 1 >= 32:
                            sram.ptr["mod_small_unsigned_end_ptr"]["i"] += 1
                            sram.ptr["mod_small_unsigned_end_ptr"]["j"] = 0
                        else:
                            sram.ptr["mod_small_unsigned_end_ptr"]["j"] += 1
                    elif (
                        color == COLORS["GREEN"]
                        and i == sram.ptr["add_mul_small_end_ptr"]["i"]
                    ):
                        if sram.ptr["add_mul_small_end_ptr"]["j"] + 1 >= 32:
                            sram.ptr["add_mul_small_end_ptr"]["i"] += 1
                            sram.ptr["add_mul_small_end_ptr"]["j"] = 0
                        else:
                            sram.ptr["add_mul_small_end_ptr"]["j"] += 1
                else:
                    write_buffer_write.append((k - 1, 1, text, color))
            write_buffer_ptr[1] -= 1

        for k, data in enumerate(mod_small_unsigned.data):

            if len(data) == 0:
                continue

            text, color = data[0]

            if text == "" or mod_small_unsigned.cnt[k] > 0:
                continue

            i = int(text.replace("x", "").split("_")[0])
            j = int(text.split("_")[1])
            if j % 2 == 0:
                if w_flag_0.flag:
                    write_buffer_write.append(
                        (write_buffer_ptr[j % 2], j % 2, text, COLORS["YELLOW"])
                    )
                    write_buffer_ptr[j % 2] += 1
                else:
                    sram_write.append((i, j, text, COLORS["YELLOW"]))
                    w_flag_0.flag = True
                    if sram.ptr["mod_small_unsigned_end_ptr"]["j"] + 1 >= 32:
                        sram.ptr["mod_small_unsigned_end_ptr"]["i"] += 1
                        sram.ptr["mod_small_unsigned_end_ptr"]["j"] = 0
                    else:
                        sram.ptr["mod_small_unsigned_end_ptr"]["j"] += 1
            else:
                if w_flag_1.flag:
                    write_buffer_write.append(
                        (write_buffer_ptr[j % 2], j % 2, text, COLORS["YELLOW"])
                    )
                    write_buffer_ptr[j % 2] += 1
                else:
                    sram_write.append((i, j, text, COLORS["YELLOW"]))
                    w_flag_1.flag = True
                    if sram.ptr["mod_small_unsigned_end_ptr"]["j"] + 1 >= 32:
                        sram.ptr["mod_small_unsigned_end_ptr"]["i"] += 1
                        sram.ptr["mod_small_unsigned_end_ptr"]["j"] = 0
                    else:
                        sram.ptr["mod_small_unsigned_end_ptr"]["j"] += 1
            mod_small_unsigned.data[k].pop(0)

        for k, data in enumerate(add_mul_small.data):
            if len(data) == 0:
                continue

            text, color = data[0]

            if text == "":
                continue

            i = int(text.replace("x", "").split("_")[0])
            j = int(text.split("_")[1])

            if add_mul_small.cnt[k] <= i + 1 and add_mul_small_valid[k]:
                d, c = add_mul_small_buffer.read(k, 0)
                add_mul_small.add(k, d, c)

        for k, data in enumerate(add_mul_small.data):

            if len(data) == 0:
                continue

            text, color = data[0]

            if text == "":
                continue

            i = int(text.replace("x", "").split("_")[0])
            j = int(text.split("_")[1])

            if add_mul_small.cnt[k] == 0:
                if j % 2 == 0:
                    if w_flag_0.flag:
                        write_buffer_write.append(
                            (write_buffer_ptr[j % 2], j % 2, text, COLORS["GREEN"])
                        )
                        write_buffer_ptr[j % 2] += 1
                    else:
                        sram_write.append((i, j, text, COLORS["GREEN"]))
                        w_flag_0.flag = True
                else:
                    if w_flag_1.flag:
                        write_buffer_write.append(
                            (write_buffer_ptr[j % 2], j % 2, text, COLORS["GREEN"])
                        )
                        write_buffer_ptr[j % 2] += 1
                    else:
                        sram_write.append((i, j, text, COLORS["GREEN"]))
                        w_flag_1.flag = True
                add_mul_small.data[k].pop(0)
            elif add_mul_small.cnt[k] <= i:
                text2, color2 = data[1]
                i2 = int(text2.replace("x", "").split("_")[0])
                j2 = int(text2.split("_")[1])
                if j2 % 2 == 0:
                    if w_flag_0.flag:
                        write_buffer_write.append(
                            (write_buffer_ptr[j2 % 2], j2 % 2, text2, COLORS["GREEN"])
                        )
                        write_buffer_ptr[j2 % 2] += 1
                    else:
                        sram_write.append((i2, j2, text2, COLORS["GREEN"]))
                        w_flag_0.flag = True
                else:
                    if w_flag_1.flag:
                        write_buffer_write.append(
                            (write_buffer_ptr[j2 % 2], j2 % 2, text2, COLORS["GREEN"])
                        )
                        write_buffer_ptr[j2 % 2] += 1
                    else:
                        sram_write.append((i2, j2, text2, COLORS["GREEN"]))
                        w_flag_1.flag = True
                add_mul_small.data[k].pop(1)

        for k in sram_write:
            sram.write(*k)
        for k in mod_small_unsigned_write:
            mod_small_unsigned.write(*k)
            if sram.ptr["mod_small_unsigned_start_ptr"]["j"] + 1 >= 32:
                sram.ptr["mod_small_unsigned_start_ptr"]["i"] += 1
                sram.ptr["mod_small_unsigned_start_ptr"]["j"] = 0
            else:
                sram.ptr["mod_small_unsigned_start_ptr"]["j"] += 1
        for k in add_mul_small_write:
            add_mul_small.write(*k)
            if sram.ptr["add_mul_small_start_ptr"]["j"] + 1 >= 32:
                sram.ptr["add_mul_small_start_ptr"]["i"] += 1
                sram.ptr["add_mul_small_start_ptr"]["j"] = 0
            else:
                sram.ptr["add_mul_small_start_ptr"]["j"] += 1
        for k in add_mul_small_buffer_write:
            add_mul_small_buffer.write(*k)
        for k in write_buffer_write:
            write_buffer.write(*k)

        # counter
        input_cnt += 1

        render(prs)


# --- PPT settings --- #
def init_ppt():
    TEMPLATE_FILE = "template.pptx"
    prs = Presentation(TEMPLATE_FILE)

    # Add title slide
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = TITLE
    subtitle.text = "Student: Chao-En Kuo\nAdvisor: Hsie-Chia Chang"

    return prs


def render(prs: pptx.presentation.Presentation):
    trace_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(trace_layout)

    for component in components:
        component.render(slide)


if __name__ == "__main__":
    prs = init_ppt()
    main_logic(prs)
    prs.save(OUTPUT_FILE)
